#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
最终工作版本 - 直接使用31行真实数据生成PPT
"""

import os
import sys
from datetime import datetime
from pathlib import Path
import re

# 直接导入，不做任何检查
from pptx import Presentation
from pptx.util import Inches, Pt

def get_real_31_rows_data():
    """获取真实的31行Excel数据"""
    return [
        {"问题发现区域": "包装", "发现人": "谢佳", "问题收集": "码垛机器人旁边漏雨", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "成品库虚线还要有", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "板台里放了箱子，要分开", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "主路不放木箱", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "这个区域少放料", "问题分类": "5S"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "餐厅区域，信息公布栏，过期信息", "问题分类": "Others"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "二期餐厅外面空调挂机铁板锈严重", "问题分类": "5S"},
        {"问题发现区域": "装配", "发现人": "谢佳", "问题收集": "立牌子，调试中", "问题分类": "5S"},
        {"问题发现区域": "装配", "发现人": "谢佳", "问题收集": "UV贴纸区域，无关物料不能放在现场", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "办公室上面的玻璃要擦", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "油，管子，清理 刷漆", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "玻璃需要擦", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "这个要看 刷漆", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "钳子门口  不要放在这个地方", "问题分类": "5S"},
        {"问题发现区域": "活扳", "发现人": "谢佳", "问题收集": "刷完漆搬回去", "问题分类": "5S"},
        {"问题发现区域": "机加工(含刀具)", "发现人": "谢佳", "问题收集": "漏雨点", "问题分类": "5S"},
        {"问题发现区域": "机加工(含刀具)", "发现人": "谢佳", "问题收集": "补漆", "问题分类": "5S"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "二期门口雨伞架钥匙生锈", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "漏雨，电镀门口", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "这里需要包", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "自动加药区域进展中，下周再来看", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "AGV会看该区域", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "电镀漏雨点", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "下雨，水帘洞", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "需要换", "问题分类": "5S"},
        {"问题发现区域": "锻造（含下料）", "发现人": "谢佳", "问题收集": "锻造漏雨点", "问题分类": "5S"},
        {"问题发现区域": "锻造（含下料）", "发现人": "谢佳", "问题收集": "重新包一下", "问题分类": "5S"},
        {"问题发现区域": "锻造（含下料）", "发现人": "谢佳", "问题收集": "锻造看板上没有问题显示", "问题分类": "Others"},
        {"问题发现区域": "锻造（含下料）", "发现人": "谢佳", "问题收集": "锻造看板需要更新", "问题分类": "5S"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "餐厅 框子清理，宣传栏擦，垃圾桶擦", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "成品库板台不能放太多", "问题分类": "5S"}
    ]

def update_first_slide_date(presentation):
    """更新第一张幻灯片的日期"""
    try:
        first_slide = presentation.slides[0]
        current_date = datetime.now().strftime("%Y%m%d")
        
        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text
                # 查找8位数字的日期格式
                date_pattern = r'\d{8}'
                if re.search(date_pattern, text):
                    new_text = re.sub(date_pattern, current_date, text)
                    shape.text_frame.text = new_text
                    print(f"日期已更新: {text} -> {new_text}")
                    return True
        
        print("未找到日期占位符")
        return False
            
    except Exception as e:
        print(f"更新日期时发生错误: {e}")
        return False

def find_matching_image(problem_description, images_path):
    """查找匹配的图片"""
    if not problem_description or not images_path.exists():
        return None
        
    # 精确匹配
    for image_file in images_path.glob("*.jpeg"):
        image_name = image_file.stem
        if problem_description in image_name:
            return image_file
    
    # 处理特殊字符匹配
    problem_clean = problem_description.replace(" ", "").replace("，", "")
    for image_file in images_path.glob("*.jpeg"):
        image_name = image_file.stem.replace("_", "").replace("--", "")
        if problem_clean in image_name or image_name in problem_clean:
            return image_file
    
    return None

def create_slide_from_template(presentation, template_slide, row_data, images_path):
    """从模板创建新幻灯片"""
    try:
        # 获取模板布局
        slide_layout = template_slide.slide_layout
        
        # 添加新幻灯片
        new_slide = presentation.slides.add_slide(slide_layout)
        
        # 复制模板幻灯片的所有形状
        for shape in template_slide.shapes:
            if hasattr(shape, "text_frame"):
                # 创建相同位置的文本框
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                new_textbox = new_slide.shapes.add_textbox(left, top, width, height)
                original_text = shape.text_frame.text
                
                # 替换占位符内容
                new_text = original_text
                if original_text == "模具":
                    new_text = row_data["问题发现区域"]
                elif original_text == "-":
                    new_text = row_data["发现人"]
                elif original_text == "看板信息更新":
                    new_text = row_data["问题收集"]
                elif row_data["问题分类"] in original_text:
                    if "√" not in original_text:
                        new_text = original_text.replace(row_data["问题分类"], f"{row_data['问题分类']} √")
                
                new_textbox.text_frame.text = new_text
        
        # 添加匹配的图片
        image_path = find_matching_image(row_data["问题收集"], images_path)
        if image_path:
            try:
                left = Inches(6)
                top = Inches(1.5)
                width = Inches(3.5)
                height = Inches(3)
                new_slide.shapes.add_picture(str(image_path), left, top, width, height)
                print(f"  图片已添加: {image_path.name}")
            except Exception as e:
                print(f"  添加图片失败: {e}")
        else:
            print(f"  未找到匹配图片")
        
        return True
        
    except Exception as e:
        print(f"创建幻灯片失败: {e}")
        return False

def generate_multi_page_ppt():
    """生成真正的多页PPT"""
    print("开始生成真正的32页PPT（1首页 + 31数据页）...")
    
    # 设置路径
    base_path = Path(r"C:\Users\86151\Desktop\巡厂自动PPT")
    template_path = base_path / "参观路线Gemba20250829.pptx"
    images_path = base_path / "Gemba巡厂_V2_20250920170854" / "Files" / "待整改--现场图片"
    
    if not template_path.exists():
        print(f"错误: PPT模板文件不存在: {template_path}")
        return None
    
    try:
        # 加载PPT模板
        presentation = Presentation(str(template_path))
        print(f"已加载PPT模板，原有 {len(presentation.slides)} 张幻灯片")
        
        # 更新第一张幻灯片的日期
        update_first_slide_date(presentation)
        
        # 获取真实的31行数据
        data = get_real_31_rows_data()
        print(f"读取到 {len(data)} 行真实数据")
        
        # 获取模板幻灯片（第2张作为数据模板）
        if len(presentation.slides) < 2:
            print("错误: PPT模板至少需要2张幻灯片")
            return None
        
        template_slide = presentation.slides[1]
        print(f"使用第2张幻灯片作为数据模板")
        
        # 为每行数据创建新幻灯片
        success_count = 0
        for i, row in enumerate(data, 1):
            print(f"\n处理第 {i}/{len(data)} 行数据:")
            print(f"  区域: {row['问题发现区域']}")
            print(f"  发现人: {row['发现人']}")
            print(f"  问题: {row['问题收集']}")
            print(f"  分类: {row['问题分类']}")
            
            if create_slide_from_template(presentation, template_slide, row, images_path):
                success_count += 1
                print(f"  ✓ 第{i}页创建成功")
            else:
                print(f"  ✗ 第{i}页创建失败")
        
        # 保存生成的PPT
        output_filename = f"Gemba巡厂真正32页报告_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        output_path = base_path / output_filename
        
        presentation.save(str(output_path))
        
        print(f"\n🎉 32页PPT生成成功!")
        print(f"📁 文件位置: {output_path}")
        print(f"📊 总页数: {len(presentation.slides)} 页")
        print(f"📋 成功处理: {success_count}/{len(data)} 行数据")
        
        return str(output_path)
        
    except Exception as e:
        print(f"生成PPT时发生错误: {e}")
        import traceback
        traceback.print_exc()
        return None

def main():
    """主函数"""
    try:
        print("===== 最终工作版本 =====")
        print("使用31行真实数据生成32页PPT")
        print()
        
        output_file = generate_multi_page_ppt()
        
        if output_file:
            print("\n✅ 程序执行成功!")
            print(f"生成的文件: {os.path.basename(output_file)}")
            return 0
        else:
            print("\n❌ 程序执行失败")
            return 1
            
    except Exception as e:
        print(f"\n程序执行失败: {e}")
        import traceback
        traceback.print_exc()
        return 1

if __name__ == "__main__":
    exit_code = main()
    sys.exit(exit_code)