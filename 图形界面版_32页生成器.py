#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
å›¾å½¢ç•Œé¢ç‰ˆ32é¡µPPTç”Ÿæˆå™¨ - Webéƒ¨ç½²å…¼å®¹ç‰ˆæœ¬
"""

# ğŸš¨ Critical: Disable ALL GUI backends to prevent libtk8.6.so error
import os
import sys
os.environ['MPLBACKEND'] = 'Agg'  # Disable matplotlib GUI backend
os.environ['DISPLAY'] = ''        # Disable X11 display
os.environ['QT_QPA_PLATFORM'] = 'offscreen'  # Disable Qt GUI
os.environ['SDL_VIDEODRIVER'] = 'dummy'      # Disable SDL video

# Disable pandas plotting backends that might trigger tkinter
import warnings
warnings.filterwarnings('ignore', category=UserWarning, module='.*')

import json
from datetime import datetime
from pathlib import Path
import re

# å¯¼å…¥åº“
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd  # ç”¨äºExcelè¯»å–

# é…ç½®æ–‡ä»¶è·¯å¾„
CONFIG_FILE = "gemba_config.json"

def load_config():
    """åŠ è½½é…ç½®æ–‡ä»¶"""
    try:
        if os.path.exists(CONFIG_FILE):
            with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception as e:
        print(f"åŠ è½½é…ç½®å¤±è´¥: {e}")
    
    # é»˜è®¤é…ç½®
    return {
        "last_ppt_folder": os.path.expanduser("~/Desktop"),
        "last_zip_folder": os.path.expanduser("~/Desktop"),
        "last_ppt_file": "",
        "last_zip_file": ""
    }

def save_config(config):
    """ä¿å­˜é…ç½®æ–‡ä»¶"""
    try:
        with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
            json.dump(config, f, ensure_ascii=False, indent=2)
        print("é…ç½®å·²ä¿å­˜")
    except Exception as e:
        print(f"ä¿å­˜é…ç½®å¤±è´¥: {e}")

def select_files_web_compatible(ppt_file=None, zip_file=None, output_folder=None):
    """Webå…¼å®¹çš„æ–‡ä»¶å¤„ç†å‡½æ•° - ä¸ä½¿ç”¨GUIå¯¹è¯æ¡†"""
    print("Webæ¨¡å¼ï¼šä½¿ç”¨å‚æ•°ä¼ é€’çš„æ–‡ä»¶è·¯å¾„")
    
    # æ£€æŸ¥å¿…éœ€çš„æ–‡ä»¶å‚æ•°
    if not ppt_file:
        print("é”™è¯¯: æœªæä¾›PPTæ¨¡æ¿æ–‡ä»¶è·¯å¾„")
        return None, None, None
        
    if not zip_file:
        print("é”™è¯¯: æœªæä¾›ZIPå‹ç¼©åŒ…æ–‡ä»¶è·¯å¾„")
        return None, None, None
        
    if not output_folder:
        output_folder = "."  # é»˜è®¤ä½¿ç”¨å½“å‰ç›®å½•
        print("ä½¿ç”¨é»˜è®¤è¾“å‡ºæ–‡ä»¶å¤¹: å½“å‰ç›®å½•")
    
    # éªŒè¯æ–‡ä»¶å­˜åœ¨æ€§
    if not os.path.exists(ppt_file):
        print(f"é”™è¯¯: PPTæ¨¡æ¿æ–‡ä»¶ä¸å­˜åœ¨: {ppt_file}")
        return None, None, None
        
    if not os.path.exists(zip_file):
        print(f"é”™è¯¯: ZIPæ–‡ä»¶ä¸å­˜åœ¨: {zip_file}")
        return None, None, None
    
    print(f"æ–‡ä»¶éªŒè¯æˆåŠŸ:")
    print(f"PPTæ¨¡æ¿: {os.path.basename(ppt_file)}")
    print(f"ZIPæ–‡ä»¶: {os.path.basename(zip_file)}")
    print(f"è¾“å‡ºä½ç½®: {output_folder}")
    
    return ppt_file, zip_file, output_folder

# ä¿æŒåŸå‡½æ•°åä½†é‡å®šå‘åˆ°æ–°çš„Webå…¼å®¹ç‰ˆæœ¬
def select_files():
    """å‘åå…¼å®¹çš„åŒ…è£…å‡½æ•°"""
    print("è­¦å‘Š: åŸGUIç‰ˆæœ¬å·²ç¦ç”¨ï¼Œè¯·ä½¿ç”¨Webç•Œé¢æˆ–ç›´æ¥è°ƒç”¨generate_ppt_with_user_files")
    return None, None, None

def read_excel_data(excel_path):
    """ä»Excelæ–‡ä»¶åŠ¨æ€è¯»å–æ•°æ®ï¼Œæ›¿ä»£ç¡¬ç¼–ç æ•°æ®"""
    try:
        print(f"æ­£åœ¨è¯»å–Excelæ–‡ä»¶: {excel_path}")
        
        # ä½¿ç”¨pandasè¯»å–Excelæ–‡ä»¶
        df = pd.read_excel(excel_path)
        print(f"Excelæ–‡ä»¶è¯»å–æˆåŠŸï¼Œå…± {len(df)} è¡Œæ•°æ®")
        
        # æ•°æ®éªŒè¯ï¼šæ£€æŸ¥å¿…éœ€çš„åˆ—æ˜¯å¦å­˜åœ¨
        required_columns = ["é—®é¢˜å‘ç°åŒºåŸŸ", "å‘ç°äºº", "é—®é¢˜æ”¶é›†", "é—®é¢˜åˆ†ç±»"]
        missing_columns = [col for col in required_columns if col not in df.columns]
        
        if missing_columns:
            print(f"è­¦å‘Šï¼šExcelæ–‡ä»¶ç¼ºå°‘å¿…éœ€åˆ—: {missing_columns}")
            # ä½¿ç”¨é»˜è®¤å€¼å¡«å……ç¼ºå¤±åˆ—
            for col in missing_columns:
                df[col] = "æœªçŸ¥"
        
        # è¿‡æ»¤ç©ºè¡Œå’Œæ— æ•ˆæ•°æ®
        df_cleaned = df.dropna(subset=["é—®é¢˜æ”¶é›†"]).copy()
        print(f"æ¸…ç†åæœ‰æ•ˆæ•°æ®: {len(df_cleaned)} è¡Œ")
        
        # è½¬æ¢ä¸ºæ ‡å‡†æ ¼å¼
        data_list = []
        for _, row in df_cleaned.iterrows():
            data_row = {
                "é—®é¢˜å‘ç°åŒºåŸŸ": str(row.get("é—®é¢˜å‘ç°åŒºåŸŸ", "æœªçŸ¥")).strip(),
                "å‘ç°äºº": str(row.get("å‘ç°äºº", "æœªçŸ¥")).strip(),
                "é—®é¢˜æ”¶é›†": str(row.get("é—®é¢˜æ”¶é›†", "")).strip(),
                "é—®é¢˜åˆ†ç±»": str(row.get("é—®é¢˜åˆ†ç±»", "Others")).strip()
            }
            # åªæ·»åŠ éç©ºçš„é—®é¢˜è®°å½•
            if data_row["é—®é¢˜æ”¶é›†"]:
                data_list.append(data_row)
        
        print(f"æœ€ç»ˆå¤„ç†æ•°æ®: {len(data_list)} è¡Œ")
        return data_list
        
    except Exception as e:
        print(f"è¯»å–Excelæ–‡ä»¶å¤±è´¥: {e}")
        print("ä½¿ç”¨å¤‡ç”¨ç¡¬ç¼–ç æ•°æ®...")
        # å‘ç”Ÿé”™è¯¯æ—¶è¿”å›åŸæœ‰çš„ç¡¬ç¼–ç æ•°æ®ä½œä¸ºå¤‡ç”¨
        return get_all_31_rows_backup()

def get_all_31_rows_backup():
    """å¤‡ç”¨ç¡¬ç¼–ç æ•°æ®å‡½æ•°ï¼ˆåŸget_all_31_rowsé‡å‘½åï¼‰"""
    return [
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "åŒ…è£…", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "ç å›æœºå™¨äººæ—è¾¹æ¼é›¨", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "æˆå“åº“ã€ç©ºæŸ„åº“", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "æˆå“åº“è™šçº¿è¿˜è¦æœ‰", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "æˆå“åº“ã€ç©ºæŸ„åº“", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "æ¿å°é‡Œæ”¾äº†ç®±å­ï¼Œè¦åˆ†å¼€", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "æˆå“åº“ã€ç©ºæŸ„åº“", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "ä¸»è·¯ä¸æ”¾æœ¨ç®±", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "æˆå“åº“ã€ç©ºæŸ„åº“", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "è¿™ä¸ªåŒºåŸŸå°‘æ”¾æ–™", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "å…¬å…±åŒºåŸŸ", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "é¤å…åŒºåŸŸï¼Œä¿¡æ¯å…¬å¸ƒæ ï¼Œè¿‡æœŸä¿¡æ¯", "é—®é¢˜åˆ†ç±»": "Others"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "å…¬å…±åŒºåŸŸ", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "äºŒæœŸé¤å…å¤–é¢ç©ºè°ƒæŒ‚æœºé“æ¿é”ˆä¸¥é‡", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "è£…é…", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "ç«‹ç‰Œå­ï¼Œè°ƒè¯•ä¸­", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "è£…é…", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "UVè´´çº¸åŒºåŸŸï¼Œæ— å…³ç‰©æ–™ä¸èƒ½æ”¾åœ¨ç°åœº", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é’³å­å†·é”»", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "åŠå…¬å®¤ä¸Šé¢çš„ç»ç’ƒè¦æ“¦", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é’³å­å†·é”»", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "æ²¹ï¼Œç®¡å­ï¼Œæ¸…ç† åˆ·æ¼†", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é’³å­å†·é”»", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "ç»ç’ƒéœ€è¦æ“¦", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é’³å­å†·é”»", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "è¿™ä¸ªè¦çœ‹ åˆ·æ¼†", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é’³å­å†·é”»", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "é’³å­é—¨å£  ä¸è¦æ”¾åœ¨è¿™ä¸ªåœ°æ–¹", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "æ´»æ‰³", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "åˆ·å®Œæ¼†æ¬å›å»", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "æœºåŠ å·¥(å«åˆ€å…·)", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "æ¼é›¨ç‚¹", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "æœºåŠ å·¥(å«åˆ€å…·)", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "è¡¥æ¼†", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "å…¬å…±åŒºåŸŸ", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "äºŒæœŸé—¨å£é›¨ä¼æ¶é’¥åŒ™ç”Ÿé”ˆ", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "ç”µé•€", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "æ¼é›¨ï¼Œç”µé•€é—¨å£", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "ç”µé•€", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "è¿™é‡Œéœ€è¦åŒ…", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "ç”µé•€", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "è‡ªåŠ¨åŠ è¯åŒºåŸŸè¿›å±•ä¸­ï¼Œä¸‹å‘¨å†æ¥çœ‹", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "ç”µé•€", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "AGVä¼šçœ‹è¯¥åŒºåŸŸ", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "ç”µé•€", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "ç”µé•€æ¼é›¨ç‚¹", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "ç”µé•€", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "ä¸‹é›¨ï¼Œæ°´å¸˜æ´", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "ç”µé•€", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "éœ€è¦æ¢", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é”»é€ ï¼ˆå«ä¸‹æ–™ï¼‰", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "é”»é€ æ¼é›¨ç‚¹", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é”»é€ ï¼ˆå«ä¸‹æ–™ï¼‰", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "é‡æ–°åŒ…ä¸€ä¸‹", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é”»é€ ï¼ˆå«ä¸‹æ–™ï¼‰", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "é”»é€ çœ‹æ¿ä¸Šæ²¡æœ‰é—®é¢˜æ˜¾ç¤º", "é—®é¢˜åˆ†ç±»": "Others"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "é”»é€ ï¼ˆå«ä¸‹æ–™ï¼‰", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "é”»é€ çœ‹æ¿éœ€è¦æ›´æ–°", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "å…¬å…±åŒºåŸŸ", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "é¤å… æ¡†å­æ¸…ç†ï¼Œå®£ä¼ æ æ“¦ï¼Œåƒåœ¾æ¡¶æ“¦", "é—®é¢˜åˆ†ç±»": "5S"},
        {"é—®é¢˜å‘ç°åŒºåŸŸ": "æˆå“åº“ã€ç©ºæŸ„åº“", "å‘ç°äºº": "è°¢ä½³", "é—®é¢˜æ”¶é›†": "æˆå“åº“æ¿å°ä¸èƒ½æ”¾å¤ªå¤š", "é—®é¢˜åˆ†ç±»": "5S"}
    ]

def get_category_mapping():
    """è·å–åˆ†ç±»æ˜ å°„"""
    return {
        "Safety": "A",
        "Quality": "B",
        "Efficiency": "C",
        "5S": "D",
        "Cost": "E",
        "Delivery": "F",
        "Others": "G"
    }

def find_matching_image(problem_description, images_path):
    """å¢å¼ºçš„å›¾ç‰‡åŒ¹é…å‡½æ•°"""
    if not problem_description or not images_path.exists():
        return None
    
    # æ–¹æ³•1: ç²¾ç¡®åŒ¹é…
    for img_file in images_path.glob("*.jpeg"):
        if problem_description in img_file.stem:
            return img_file
    
    # æ–¹æ³•2: åå‘åŒ¹é…
    for img_file in images_path.glob("*.jpeg"):
        img_name = img_file.stem
        if img_name in problem_description:
            return img_file
    
    # æ–¹æ³•3: æ¸…ç†ç‰¹æ®Šå­—ç¬¦ååŒ¹é…
    problem_clean = problem_description.replace(" ", "").replace("ï¼Œ", "").replace("ã€‚", "")
    for img_file in images_path.glob("*.jpeg"):
        img_clean = img_file.stem.replace("_", "").replace("-", "").replace("--", "")
        if problem_clean in img_clean or img_clean in problem_clean:
            return img_file
    
    # æ–¹æ³•4: å…³é”®è¯åŒ¹é…
    keywords = problem_description.split()
    for img_file in images_path.glob("*.jpeg"):
        img_name = img_file.stem
        for keyword in keywords:
            if len(keyword) > 1 and keyword in img_name:
                return img_file
    
    return None

def handle_circle_markers(slide, target_category):
    """å¤„ç†åœ†å½¢æ ‡è®° A-G ç³»ç»Ÿ"""
    category_mapping = get_category_mapping()
    target_letter = category_mapping.get(target_category)
    
    if not target_letter:
        print(f"    è­¦å‘Š: æœªçŸ¥åˆ†ç±»: {target_category}")
        return
    
    print(f"    å¤„ç†åœ†å½¢æ ‡è®°: {target_category} -> {target_letter}")
    
    # æŸ¥æ‰¾æ‰€æœ‰åœ†å½¢å’Œæ–‡æœ¬å½¢çŠ¶
    circles_to_remove = []
    target_circle = None
    
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            text = shape.text_frame.text.strip()
            
            # æ£€æŸ¥æ˜¯å¦æ˜¯åˆ†ç±»å­—æ¯æ ‡è®°
            if text in ["A", "B", "C", "D", "E", "F", "G"]:
                if text == target_letter:
                    # è¿™æ˜¯ç›®æ ‡åœ†åœˆï¼Œæ·»åŠ å‹¾é€‰æ ‡è®°
                    shape.text_frame.text = "V"
                    target_circle = shape
                    print(f"      [V] åœ¨åœ†åœˆ {text} ä¸­æ·»åŠ å‹¾é€‰")
                else:
                    # è¿™æ˜¯å…¶ä»–åœ†åœˆï¼Œæ ‡è®°ä¸ºåˆ é™¤
                    circles_to_remove.append(shape)
                    print(f"      [X] æ ‡è®°åˆ é™¤åœ†åœˆ {text}")
    
    # åˆ é™¤æœªæ ‡è®°çš„åœ†åœˆ
    for shape in circles_to_remove:
        try:
            # åˆ é™¤å½¢çŠ¶çš„æ–¹æ³•
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"      å·²åˆ é™¤æœªä½¿ç”¨çš„åœ†åœˆ")
        except Exception as e:
            print(f"      åˆ é™¤åœ†åœˆå¤±è´¥: {e}")
    
    if target_circle:
        print(f"    [V] åœ†åœˆæ ‡è®°å¤„ç†å®Œæˆ: {target_category}")
    else:
        print(f"    [!] æœªæ‰¾åˆ°ç›®æ ‡åœ†åœˆ: {target_letter}")

def extract_zip_and_find_files(zip_path):
    """è§£å‹ZIPæ–‡ä»¶å¹¶æŸ¥æ‰¾Excelå’Œå›¾ç‰‡"""
    import zipfile
    import tempfile
    
    try:
        zip_path = Path(zip_path)
        temp_dir = Path(tempfile.mkdtemp())
        
        # è§£å‹ZIPæ–‡ä»¶
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        print(f"ZIPæ–‡ä»¶å·²è§£å‹åˆ°: {temp_dir}")
        
        # æŸ¥æ‰¾Excelæ–‡ä»¶
        excel_files = list(temp_dir.rglob("*.xlsx"))
        if not excel_files:
            raise FileNotFoundError("æœªåœ¨ZIPæ–‡ä»¶ä¸­æ‰¾åˆ°Excelæ–‡ä»¶")
        
        excel_path = excel_files[0]
        print(f"æ‰¾åˆ°Excelæ–‡ä»¶: {excel_path}")
        
        # æŸ¥æ‰¾å›¾ç‰‡æ–‡ä»¶å¤¹
        image_folders = []
        for folder in temp_dir.rglob("*"):
            if folder.is_dir() and ("å›¾ç‰‡" in folder.name or "ç…§ç‰‡" in folder.name):
                image_folders.append(folder)
        
        if not image_folders:
            # å¦‚æœæ²¡æœ‰æ‰¾åˆ°ä¸“é—¨çš„å›¾ç‰‡æ–‡ä»¶å¤¹ï¼ŒæŸ¥æ‰¾åŒ…å«jpegæ–‡ä»¶çš„æ–‡ä»¶å¤¹
            for folder in temp_dir.rglob("*"):
                if folder.is_dir() and list(folder.glob("*.jpeg")):
                    image_folders.append(folder)
        
        if not image_folders:
            raise FileNotFoundError("æœªåœ¨ZIPæ–‡ä»¶ä¸­æ‰¾åˆ°å›¾ç‰‡æ–‡ä»¶å¤¹")
        
        images_path = image_folders[0]
        print(f"æ‰¾åˆ°å›¾ç‰‡æ–‡ä»¶å¤¹: {images_path}")
        
        return excel_path, images_path, temp_dir
        
    except Exception as e:
        print(f"è§£å‹ZIPæ–‡ä»¶å¤±è´¥: {e}")
        return None, None, None

def generate_ppt_with_user_files(ppt_file, zip_file, output_folder):
    """ä½¿ç”¨ç”¨æˆ·é€‰æ‹©çš„æ–‡ä»¶ç”ŸæˆPPT"""
    print("å¼€å§‹ç”ŸæˆPPT...")
    
    try:
        # è§£å‹ZIPæ–‡ä»¶å¹¶æŸ¥æ‰¾ç›¸å…³æ–‡ä»¶
        excel_path, images_path, temp_dir = extract_zip_and_find_files(zip_file)
        
        if not excel_path or not images_path:
            print("æ— æ³•æ‰¾åˆ°Excelæ–‡ä»¶æˆ–å›¾ç‰‡æ–‡ä»¶å¤¹")
            return None
        
        # æ˜¾ç¤ºæ‰¾åˆ°çš„å›¾ç‰‡æ•°é‡
        images = list(images_path.glob("*.jpeg"))
        print(f"å‘ç° {len(images)} å¼ å›¾ç‰‡")
        
        # åŠ è½½PPTæ¨¡æ¿
        prs = Presentation(ppt_file)
        print(f"åŠ è½½PPTæ¨¡æ¿æˆåŠŸï¼ŒåŸæœ‰ {len(prs.slides)} å¼ å¹»ç¯ç‰‡")
        
        # æ›´æ–°ç¬¬ä¸€é¡µæ—¥æœŸ
        first_slide = prs.slides[0]
        current_date = datetime.now().strftime("%Y%m%d")
        
        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text
                if re.search(r'\d{8}', text):
                    new_text = re.sub(r'\d{8}', current_date, text)
                    shape.text_frame.text = new_text
                    print(f"æ—¥æœŸå·²æ›´æ–°: {text} -> {new_text}")
                    break
        
        # è·å–Excelæ•°æ®ï¼ˆåŠ¨æ€è¡Œæ•°ï¼‰
        data = read_excel_data(excel_path)
        print(f"ä»Excelè¯»å–åˆ° {len(data)} è¡Œæ•°æ®ï¼Œå‡†å¤‡å¤„ç†")
        
        # è·å–æ¨¡æ¿å¹»ç¯ç‰‡
        if len(prs.slides) < 2:
            print("é”™è¯¯: PPTæ¨¡æ¿éœ€è¦è‡³å°‘2å¼ å¹»ç¯ç‰‡")
            return None
            
        template_slide = prs.slides[1]
        
        # ä¸ºæ¯è¡Œæ•°æ®åˆ›å»ºå¹»ç¯ç‰‡
        created_count = 0
        images_found = 0
        
        for i, row in enumerate(data, 1):
            print(f"\nåˆ›å»ºç¬¬ {i+1} é¡µ: {row['é—®é¢˜æ”¶é›†'][:30]}...")
            
            try:
                # æ·»åŠ æ–°å¹»ç¯ç‰‡
                slide_layout = template_slide.slide_layout
                new_slide = prs.slides.add_slide(slide_layout)
                
                # å¤åˆ¶æ¨¡æ¿å†…å®¹
                for shape in template_slide.shapes:
                    if hasattr(shape, "text_frame"):
                        # åˆ›å»ºæ–°æ–‡æœ¬æ¡†
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        new_textbox = new_slide.shapes.add_textbox(left, top, width, height)
                        original_text = shape.text_frame.text
                        
                        # æ›¿æ¢å ä½ç¬¦å†…å®¹
                        if original_text == "æ¨¡å…·":
                            new_textbox.text_frame.text = row["é—®é¢˜å‘ç°åŒºåŸŸ"]
                            print(f"    å ä½ç¬¦1: æ¨¡å…· -> {row['é—®é¢˜å‘ç°åŒºåŸŸ']}")
                        elif original_text == "-":
                            new_textbox.text_frame.text = row["å‘ç°äºº"]
                            print(f"    å ä½ç¬¦2: - -> {row['å‘ç°äºº']}")
                        elif "çœ‹æ¿ä¿¡æ¯æ›´æ–°" in original_text or original_text == "çœ‹æ¿ä¿¡æ¯æ›´æ–°":
                            new_textbox.text_frame.text = row["é—®é¢˜æ”¶é›†"]
                            print(f"    å ä½ç¬¦4: {original_text} -> {row['é—®é¢˜æ”¶é›†']}")
                        else:
                            new_textbox.text_frame.text = original_text
                
                # å¤„ç†åœ†å½¢æ ‡è®°ç³»ç»Ÿ
                handle_circle_markers(new_slide, row["é—®é¢˜åˆ†ç±»"])
                
                # æ·»åŠ å›¾ç‰‡åˆ°å·¦è¾¹ - å®Œç¾ä½ç½®
                image_path = find_matching_image(row["é—®é¢˜æ”¶é›†"], images_path)
                if image_path:
                    try:
                        left = Inches(0.5)  # å·¦è¾¹ä½ç½®
                        top = Inches(2.1)   # å®Œç¾é«˜åº¦
                        width = Inches(3.5)
                        height = Inches(2.8)  # å®Œç¾é«˜åº¦
                        new_slide.shapes.add_picture(str(image_path), left, top, width, height)
                        print(f"    [V] å›¾ç‰‡å·²æ·»åŠ : {image_path.name}")
                        images_found += 1
                    except Exception as e:
                        print(f"    [X] å›¾ç‰‡æ·»åŠ å¤±è´¥: {e}")
                else:
                    print(f"    [X] æœªæ‰¾åˆ°åŒ¹é…å›¾ç‰‡")
                
                created_count += 1
                
            except Exception as e:
                print(f"  åˆ›å»ºç¬¬{i+1}é¡µå¤±è´¥: {e}")
        
        # åˆ é™¤åŸå§‹çš„ç¬¬äºŒé¡µæ¨¡æ¿å¹»ç¯ç‰‡
        if len(prs.slides) > 1:
            try:
                # åˆ é™¤ç¬¬äºŒå¼ å¹»ç¯ç‰‡ï¼ˆæ¨¡æ¿é¡µï¼‰
                slide_to_remove = prs.slides[1]
                slide_id = slide_to_remove.slide_id
                
                # ä»å¹»ç¯ç‰‡åˆ—è¡¨ä¸­ç§»é™¤
                for slide_rel in list(prs.slides._sldIdLst):
                    if slide_rel.id == slide_id:
                        prs.slides._sldIdLst.remove(slide_rel)
                        print("[V] å·²åˆ é™¤åŸå§‹ç¬¬äºŒé¡µæ¨¡æ¿å¹»ç¯ç‰‡")
                        break
                        
            except Exception as e:
                print(f"åˆ é™¤æ¨¡æ¿å¹»ç¯ç‰‡æ—¶å‘ç”Ÿé”™è¯¯: {e}")
        
        # ä¿å­˜PPT - ä½¿ç”¨ç®€åŒ–çš„æ–‡ä»¶å
        current_date_str = datetime.now().strftime("%Y%m%d")
        output_name = f"Gembaå·¡å‚æŠ¥å‘Š{current_date_str}.pptx"
        output_path = Path(output_folder) / output_name
        
        prs.save(str(output_path))
        
        # æ¸…ç†ä¸´æ—¶æ–‡ä»¶
        import shutil
        if temp_dir.exists():
            shutil.rmtree(temp_dir)
        
        print(f"\n[æˆåŠŸ] åŠ¨æ€PPTç”ŸæˆæˆåŠŸ!")
        print(f"æ–‡ä»¶: {output_name}")
        print(f"ä¿å­˜ä½ç½®: {output_folder}")
        print(f"æ€»é¡µæ•°: {len(prs.slides)} é¡µ")
        print(f"Excelæ•°æ®è¡Œæ•°: {len(data)} è¡Œ")
        print(f"[V] æˆåŠŸåˆ›å»ºæ•°æ®é¡µ: {created_count}/{len(data)} é¡µ")
        print(f"æˆåŠŸæ·»åŠ å›¾ç‰‡: {images_found}/{len(data)} é¡µ")
        print(f"åŸå§‹æ¨¡æ¿é¡µ: å·²åˆ é™¤")
        
        # æ˜¾ç¤ºå®Œæˆä¿¡æ¯ï¼ˆWebå…¼å®¹ç‰ˆæœ¬ï¼‰
        print("\n=== PPTç”Ÿæˆå®Œæˆ ===")
        print(f"æ–‡ä»¶å: {output_name}")
        print(f"ä¿å­˜ä½ç½®: {output_folder}")
        print(f"æ€»é¡µæ•°: {len(prs.slides)} é¡µ")
        print(f"Excelæ•°æ®: {len(data)} è¡Œ")
        print(f"æ•°æ®é¡µ: {created_count} é¡µ")
        print(f"å›¾ç‰‡: {images_found} å¼ ")
        print("==================")
        
        return str(output_path)
        
    except Exception as e:
        print(f"ç”ŸæˆPPTæ—¶å‘ç”Ÿé”™è¯¯: {e}")
        import traceback
        traceback.print_exc()
        
        # æ˜¾ç¤ºé”™è¯¯ä¿¡æ¯ï¼ˆWebå…¼å®¹ç‰ˆæœ¬ï¼‰
        print(f"\n=== PPTç”Ÿæˆå¤±è´¥ ===")
        print(f"é”™è¯¯ä¿¡æ¯: {e}")
        print("==================")
        
        return None

def main():
    """ä¸»å‡½æ•° - å¸¦å›¾å½¢ç•Œé¢çš„PPTç”Ÿæˆå™¨"""
    try:
        print("===== Gembaå·¡å‚PPTç”Ÿæˆå™¨ (å›¾å½¢ç•Œé¢ç‰ˆ) =====")
        print("å³å°†æ‰“å¼€æ–‡ä»¶é€‰æ‹©å¯¹è¯æ¡†...")
        
        # é€‰æ‹©æ–‡ä»¶
        ppt_file, zip_file, output_folder = select_files()
        
        if not all([ppt_file, zip_file, output_folder]):
            print("ç”¨æˆ·å–æ¶ˆäº†æ–‡ä»¶é€‰æ‹©")
            return 1
        
        print(f"\nç”¨æˆ·é€‰æ‹©:")
        print(f"PPTæ¨¡æ¿: {ppt_file}")
        print(f"ZIPæ–‡ä»¶: {zip_file}")
        print(f"è¾“å‡ºä½ç½®: {output_folder}")
        
        # ç”ŸæˆPPT
        output_file = generate_ppt_with_user_files(ppt_file, zip_file, output_folder)
        
        if output_file:
            print("\n[æˆåŠŸ] ç¨‹åºæ‰§è¡ŒæˆåŠŸ!")
            return 0
        else:
            print("\n[å¤±è´¥] ç¨‹åºæ‰§è¡Œå¤±è´¥")
            return 1
            
    except Exception as e:
        print(f"\nç¨‹åºæ‰§è¡Œå¤±è´¥: {e}")
        import traceback
        traceback.print_exc()
        return 1

if __name__ == "__main__":
    exit_code = main()
    sys.exit(exit_code)