#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Streamlit 版本 - Gemba巡厂PPT生成器
将 Tkinter 桌面应用转换为 Web 应用
"""

# 🚨 Critical: Disable ALL GUI backends before any imports
import os
os.environ['MPLBACKEND'] = 'Agg'  # Disable matplotlib GUI backend
os.environ['DISPLAY'] = ''        # Disable X11 display
os.environ['QT_QPA_PLATFORM'] = 'offscreen'  # Disable Qt GUI
os.environ['SDL_VIDEODRIVER'] = 'dummy'      # Disable SDL video

# Disable pandas plotting backends that might trigger tkinter
import warnings
warnings.filterwarnings('ignore', category=UserWarning, module='.*')

import streamlit as st
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
import re
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Streamlit 配置将在main()函数中初始化

def read_excel_data(excel_path):
    """从Excel文件动态读取数据，替代硬编码数据"""
    try:
        st.info(f"正在读取Excel文件: {excel_path}")
        
        # 使用pandas读取Excel文件
        df = pd.read_excel(excel_path)
        st.success(f"Excel文件读取成功，共 {len(df)} 行数据")
        
        # 数据验证：检查必需的列是否存在
        required_columns = ["问题发现区域", "发现人", "问题收集", "问题分类"]
        missing_columns = [col for col in required_columns if col not in df.columns]
        
        if missing_columns:
            st.warning(f"Excel文件缺少必需列: {missing_columns}")
            # 使用默认值填充缺失列
            for col in missing_columns:
                df[col] = "未知"
        
        # 过滤空行和无效数据
        df_cleaned = df.dropna(subset=["问题收集"]).copy()
        st.info(f"清理后有效数据: {len(df_cleaned)} 行")
        
        # 转换为标准格式
        data_list = []
        for _, row in df_cleaned.iterrows():
            data_row = {
                "问题发现区域": str(row.get("问题发现区域", "未知")).strip(),
                "发现人": str(row.get("发现人", "未知")).strip(),
                "问题收集": str(row.get("问题收集", "")).strip(),
                "问题分类": str(row.get("问题分类", "Others")).strip()
            }
            # 只添加非空的问题记录
            if data_row["问题收集"]:
                data_list.append(data_row)
        
        st.success(f"最终处理数据: {len(data_list)} 行")
        return data_list
        
    except Exception as e:
        st.error(f"读取Excel文件失败: {e}")
        st.info("使用备用硬编码数据...")
        # 发生错误时返回原有的硬编码数据作为备用
        return get_all_31_rows_backup()

def get_all_31_rows_backup():
    """备用硬编码数据函数"""
    return [
        {"问题发现区域": "包装", "发现人": "谢佳", "问题收集": "码垛机器人旁边漏雨", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "成品库虚线还要有", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "板台里放了箱子，要分开", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "主路不放木箱", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "这个区域少放料", "问题分类": "5S"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "餐厅区域，信息公布栏，过期信息", "问题分类": "Others"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "二期餐厅外面空调挂机铁板锈严重", "问题分类": "5S"},
        {"问题发现区域": "装配", "发现人": "谢佳", "问题收集": "立牌子，调试中", "问题分类": "5S"},
        {"问题发现区域": "装配", "发现人": "谢佳", "问题收集": "UV贴纸区域，无关物料不能放在现场", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "办公室上面的玻璃要擦", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "油，管子，清理 刷漆", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "玻璃需要擦", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "这个要看 刷漆", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "钳子门口  不要放在这个地方", "问题分类": "5S"},
        {"问题发现区域": "活扳", "发现人": "谢佳", "问题收集": "刷完漆搬回去", "问题分类": "5S"},
        {"问题发现区域": "机加工(含刀具)", "发现人": "谢佳", "问题收集": "漏雨点", "问题分类": "5S"},
        {"问题发现区域": "机加工(含刀具)", "发现人": "谢佳", "问题收集": "补漆", "问题分类": "5S"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "二期门口雨伞架钥匙生锈", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "漏雨，电镀门口", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "这里需要包", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "自动加药区域进展中，下周再来看", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "AGV会看该区域", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "电镀漏雨点", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "下雨，水帘洞", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "需要换", "问题分类": "5S"},
        {"问题发现区域": "锻造（含下料）", "发现人": "谢佳", "问题收集": "锻造漏雨点", "问题分类": "5S"},
        {"问题发现区域": "锻造（含下料）", "发现人": "谢佳", "问题收集": "重新包一下", "问题分类": "5S"},
        {"问题发现区域": "锻造（含下料）", "发现人": "谢佳", "问题收集": "锻造看板上没有问题显示", "问题分类": "Others"},
        {"问题发现区域": "锻造（含下料）", "发现人": "谢佳", "问题收集": "锻造看板需要更新", "问题分类": "5S"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "餐厅 框子清理，宣传栏擦，垃圾桶擦", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "成品库板台不能放太多", "问题分类": "5S"}
    ]

def get_category_mapping():
    """获取分类映射"""
    return {
        "Safety": "A",
        "Quality": "B", 
        "Efficiency": "C",
        "5S": "D",
        "Cost": "E",
        "Delivery": "F",
        "Others": "G"
    }

def find_matching_image(problem_description, images_path):
    """增强的图片匹配函数"""
    if not problem_description or not images_path.exists():
        return None
    
    # 方法1: 精确匹配
    for img_file in images_path.glob("*.jpeg"):
        if problem_description in img_file.stem:
            return img_file
    
    # 方法2: 反向匹配
    for img_file in images_path.glob("*.jpeg"):
        img_name = img_file.stem
        if img_name in problem_description:
            return img_file
    
    # 方法3: 清理特殊字符后匹配
    problem_clean = problem_description.replace(" ", "").replace("，", "").replace("。", "")
    for img_file in images_path.glob("*.jpeg"):
        img_clean = img_file.stem.replace("_", "").replace("-", "").replace("--", "")
        if problem_clean in img_clean or img_clean in problem_clean:
            return img_file
    
    # 方法4: 关键词匹配
    keywords = problem_description.split()
    for img_file in images_path.glob("*.jpeg"):
        img_name = img_file.stem
        for keyword in keywords:
            if len(keyword) > 1 and keyword in img_name:
                return img_file
    
    return None

def handle_circle_markers(slide, target_category):
    """处理圆形标记 A-G 系统"""
    category_mapping = get_category_mapping()
    target_letter = category_mapping.get(target_category)
    
    if not target_letter:
        st.warning(f"未知分类: {target_category}")
        return
    
    # 查找所有圆形和文本形状
    circles_to_remove = []
    target_circle = None
    
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            text = shape.text_frame.text.strip()
            
            # 检查是否是分类字母标记
            if text in ["A", "B", "C", "D", "E", "F", "G"]:
                if text == target_letter:
                    # 这是目标圆圈，添加勾选标记
                    shape.text_frame.text = "V"
                    target_circle = shape
                else:
                    # 这是其他圆圈，标记为删除
                    circles_to_remove.append(shape)
    
    # 删除未标记的圆圈
    for shape in circles_to_remove:
        try:
            # 删除形状的方法
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception as e:
            st.warning(f"删除圆圈失败: {e}")
    
    return target_circle is not None

def extract_zip_and_find_files(zip_path):
    """解压ZIP文件并查找Excel和图片"""
    try:
        zip_path = Path(zip_path)
        temp_dir = Path(tempfile.mkdtemp())
        
        # 解压ZIP文件
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        st.info(f"ZIP文件已解压到: {temp_dir}")
        
        # 查找Excel文件
        excel_files = list(temp_dir.rglob("*.xlsx"))
        if not excel_files:
            raise FileNotFoundError("未在ZIP文件中找到Excel文件")
        
        excel_path = excel_files[0]
        st.success(f"找到Excel文件: {excel_path.name}")
        
        # 查找图片文件夹
        image_folders = []
        for folder in temp_dir.rglob("*"):
            if folder.is_dir() and ("图片" in folder.name or "照片" in folder.name):
                image_folders.append(folder)
        
        if not image_folders:
            # 如果没有找到专门的图片文件夹，查找包含jpeg文件的文件夹
            for folder in temp_dir.rglob("*"):
                if folder.is_dir() and list(folder.glob("*.jpeg")):
                    image_folders.append(folder)
        
        if not image_folders:
            raise FileNotFoundError("未在ZIP文件中找到图片文件夹")
        
        images_path = image_folders[0]
        st.success(f"找到图片文件夹: {images_path.name}")
        
        return excel_path, images_path, temp_dir
        
    except Exception as e:
        st.error(f"解压ZIP文件失败: {e}")
        return None, None, None

def generate_ppt_streamlit(ppt_file, zip_file, output_filename):
    """Streamlit版本的PPT生成函数"""
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    try:
        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # 保存上传的文件到临时目录
            status_text.text("保存上传文件...")
            progress_bar.progress(10)
            
            ppt_path = temp_path / "template.pptx"
            zip_path = temp_path / "data.zip"
            
            with open(ppt_path, "wb") as f:
                f.write(ppt_file.getvalue())
            
            with open(zip_path, "wb") as f:
                f.write(zip_file.getvalue())
            
            # 解压ZIP文件并查找相关文件
            status_text.text("解压ZIP文件...")
            progress_bar.progress(20)
            
            excel_path, images_path, extract_dir = extract_zip_and_find_files(zip_path)
            
            if not excel_path or not images_path:
                st.error("无法找到Excel文件或图片文件夹")
                return None
            
            # 显示找到的图片数量
            images = list(images_path.glob("*.jpeg"))
            st.info(f"发现 {len(images)} 张图片")
            
            # 加载PPT模板
            status_text.text("加载PPT模板...")
            progress_bar.progress(30)
            
            prs = Presentation(ppt_path)
            st.success(f"加载PPT模板成功，原有 {len(prs.slides)} 张幻灯片")
            
            # 更新第一页日期
            status_text.text("更新日期...")
            progress_bar.progress(40)
            
            first_slide = prs.slides[0]
            current_date = datetime.now().strftime("%Y%m%d")
            
            for shape in first_slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    if re.search(r'\d{8}', text):
                        new_text = re.sub(r'\d{8}', current_date, text)
                        shape.text_frame.text = new_text
                        st.info(f"日期已更新: {text} -> {new_text}")
                        break
            
            # 获取Excel数据
            status_text.text("读取Excel数据...")
            progress_bar.progress(50)
            
            data = read_excel_data(excel_path)
            st.success(f"从Excel读取到 {len(data)} 行数据，准备处理")
            
            # 获取模板幻灯片
            if len(prs.slides) < 2:
                st.error("PPT模板需要至少2张幻灯片")
                return None
                
            template_slide = prs.slides[1]
            
            # 为每行数据创建幻灯片
            status_text.text("生成PPT页面...")
            created_count = 0
            images_found = 0
            
            for i, row in enumerate(data, 1):
                progress_value = 50 + int((i / len(data)) * 40)
                progress_bar.progress(progress_value)
                status_text.text(f"创建第 {i+1} 页: {row['问题收集'][:30]}...")
                
                try:
                    # 添加新幻灯片
                    slide_layout = template_slide.slide_layout
                    new_slide = prs.slides.add_slide(slide_layout)
                    
                    # 复制模板内容
                    for shape in template_slide.shapes:
                        if hasattr(shape, "text_frame"):
                            # 创建新文本框
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height
                            
                            new_textbox = new_slide.shapes.add_textbox(left, top, width, height)
                            original_text = shape.text_frame.text
                            
                            # 替换占位符内容
                            if original_text == "模具":
                                new_textbox.text_frame.text = row["问题发现区域"]
                            elif original_text == "-":
                                new_textbox.text_frame.text = row["发现人"]
                            elif "看板信息更新" in original_text or original_text == "看板信息更新":
                                new_textbox.text_frame.text = row["问题收集"]
                            else:
                                new_textbox.text_frame.text = original_text
                    
                    # 处理圆形标记系统
                    handle_circle_markers(new_slide, row["问题分类"])
                    
                    # 添加图片到左边
                    image_path = find_matching_image(row["问题收集"], images_path)
                    if image_path:
                        try:
                            left = Inches(0.5)
                            top = Inches(2.1)
                            width = Inches(3.5)
                            height = Inches(2.8)
                            new_slide.shapes.add_picture(str(image_path), left, top, width, height)
                            images_found += 1
                        except Exception as e:
                            st.warning(f"图片添加失败: {e}")
                    
                    created_count += 1
                    
                except Exception as e:
                    st.error(f"创建第{i+1}页失败: {e}")
            
            # 删除原始的第二页模板幻灯片
            status_text.text("清理模板页...")
            progress_bar.progress(90)
            
            if len(prs.slides) > 1:
                try:
                    slide_to_remove = prs.slides[1]
                    slide_id = slide_to_remove.slide_id
                    
                    for slide_rel in list(prs.slides._sldIdLst):
                        if slide_rel.id == slide_id:
                            prs.slides._sldIdLst.remove(slide_rel)
                            break
                            
                except Exception as e:
                    st.warning(f"删除模板幻灯片时发生错误: {e}")
            
            # 保存PPT
            status_text.text("保存PPT文件...")
            progress_bar.progress(95)
            
            output_path = temp_path / output_filename
            prs.save(str(output_path))
            
            # 读取生成的文件用于下载
            with open(output_path, "rb") as f:
                ppt_data = f.read()
            
            progress_bar.progress(100)
            status_text.text("PPT生成完成！")
            
            # 显示成功信息
            st.success("🎉 PPT生成成功!")
            
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("总页数", f"{len(prs.slides)} 页")
            with col2:
                st.metric("数据页", f"{created_count} 页")
            with col3:
                st.metric("图片", f"{images_found} 张")
            
            return ppt_data
            
    except Exception as e:
        st.error(f"生成PPT时发生错误: {e}")
        import traceback
        st.code(traceback.format_exc())
        return None

def main():
    """主函数 - Streamlit Web应用"""
    
    # 页面配置
    st.set_page_config(
        page_title="Gemba巡厂PPT生成器",
        page_icon="🏭",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # 样式设置
    st.markdown("""
    <style>
    .main-header {
        background: linear-gradient(90deg, #1f77b4, #ff7f0e);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        font-size: 3rem;
        font-weight: bold;
        text-align: center;
        margin-bottom: 2rem;
    }
    .upload-section {
        background-color: #f0f2f6;
        padding: 1rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # 添加调试信息
    st.write("🚀 Streamlit 应用已启动")
    
    # 标题
    st.markdown('<h1 class="main-header">🏭 Gemba巡厂PPT生成器</h1>', unsafe_allow_html=True)
    
    # 说明文档
    with st.expander("📖 使用说明", expanded=False):
        st.markdown("""
        ### 功能介绍
        - 🎯 **自动生成**: 基于Excel数据和图片自动生成PPT报告
        - 📊 **数据驱动**: 支持动态行数的Excel数据导入
        - 🖼️ **图片匹配**: 智能匹配问题描述与相关图片
        - 🎨 **模板系统**: 使用PPT模板确保格式统一
        
        ### 使用步骤
        1. 上传PPT模板文件 (.pptx格式)
        2. 上传包含Excel数据和图片的ZIP压缩包
        3. 设置输出文件名
        4. 点击生成按钮
        5. 下载生成的PPT文件
        
        ### 文件要求
        - **PPT模板**: 至少包含2张幻灯片（封面+模板页）
        - **ZIP文件**: 包含Excel数据文件和图片文件夹
        - **Excel格式**: 需包含"问题发现区域"、"发现人"、"问题收集"、"问题分类"列
        """)
    
    # 文件上传区域
    st.markdown('<div class="upload-section">', unsafe_allow_html=True)
    st.subheader("📁 文件上传")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("**PPT模板文件**")
        ppt_file = st.file_uploader(
            "选择PPT模板文件",
            type=['pptx'],
            help="上传包含模板格式的PowerPoint文件",
            key="ppt_upload"
        )
        
        if ppt_file:
            st.success(f"✅ 已选择: {ppt_file.name}")
            st.info(f"文件大小: {ppt_file.size / 1024:.1f} KB")
    
    with col2:
        st.markdown("**数据压缩包**")
        zip_file = st.file_uploader(
            "选择ZIP压缩包",
            type=['zip'],
            help="上传包含Excel数据和图片的ZIP文件",
            key="zip_upload"
        )
        
        if zip_file:
            st.success(f"✅ 已选择: {zip_file.name}")
            st.info(f"文件大小: {zip_file.size / 1024:.1f} KB")
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # 输出设置
    st.subheader("⚙️ 输出设置")
    
    current_date = datetime.now().strftime("%Y%m%d")
    default_filename = f"Gemba巡厂报告{current_date}.pptx"
    
    col1, col2 = st.columns([2, 1])
    with col1:
        output_filename = st.text_input(
            "输出文件名",
            value=default_filename,
            help="设置生成的PPT文件名"
        )
    
    with col2:
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🔄 重置文件名"):
            st.session_state.clear()
            st.rerun()
    
    # 生成按钮
    st.markdown("---")
    
    if st.button("🚀 生成PPT", type="primary", use_container_width=True):
        if ppt_file and zip_file:
            if not output_filename.endswith('.pptx'):
                output_filename += '.pptx'
            
            st.markdown("### 🔄 生成进度")
            
            # 生成PPT
            ppt_data = generate_ppt_streamlit(ppt_file, zip_file, output_filename)
            
            if ppt_data:
                # 提供下载按钮
                st.markdown("### 📥 下载文件")
                st.download_button(
                    label="📥 下载生成的PPT",
                    data=ppt_data,
                    file_name=output_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
                # 显示成功消息
                st.balloons()
            else:
                st.error("❌ PPT生成失败，请检查文件格式和内容")
        else:
            st.error("⚠️ 请上传所有必需文件")
            if not ppt_file:
                st.error("• 缺少PPT模板文件")
            if not zip_file:
                st.error("• 缺少ZIP数据文件")
    
    # 侧边栏信息
    with st.sidebar:
        st.header("📊 应用信息")
        st.info("""
        **版本**: v2.0 (Streamlit版)
        **更新**: 2024年
        **开发**: 巡厂自动化团队
        """)
        
        st.header("📞 技术支持")
        st.markdown("""
        如遇问题请联系技术支持:
        - 📧 Email: support@example.com
        - 📱 电话: 400-000-0000
        """)
        
        st.header("🔗 相关链接")
        st.markdown("""
        - [使用手册](https://example.com/manual)
        - [视频教程](https://example.com/tutorial)
        - [问题反馈](https://example.com/feedback)
        """)

if __name__ == "__main__":
    main()