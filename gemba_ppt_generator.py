#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Gemba巡厂PPT自动生成器
自动处理PPT模板和Excel数据，生成巡检报告PPT

功能：
1. 读取PPT模板并更新第一页日期
2. 读取Excel数据，为每行创建新的PPT页面
3. 填充占位符数据（问题发现区域、发现人、问题收集）
4. 根据问题分类在对应选项上打勾
5. 匹配并插入对应图片
6. 保存生成的PPT文件
"""

import os
import sys
import pandas as pd
import logging
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import shutil

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('gemba_ppt_generator.log', encoding='utf-8'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)

class GembaPPTGenerator:
    """Gemba巡厂PPT生成器"""
    
    def __init__(self, base_path):
        """
        初始化生成器
        
        Args:
            base_path (str): 基础路径，包含PPT模板和数据文件
        """
        self.base_path = Path(base_path)
        self.template_path = self.base_path / "参观路线Gemba20250829.pptx"
        self.excel_path = self.base_path / "Gemba巡厂_V2_20250920170854" / "Gemba巡厂_V2_20250920170854.xlsx"
        self.images_path = self.base_path / "Gemba巡厂_V2_20250920170854" / "Files" / "待整改--现场图片"
        
        # 问题分类选项映射
        self.category_options = [
            "Safety", "Efficiency", "Cost", "Quality", 
            "Delivery", "5S", "Others"
        ]
        
        # 验证路径
        self._validate_paths()
        
    def _validate_paths(self):
        """验证所有必需的文件和路径是否存在"""
        logger.info("验证文件路径...")
        
        if not self.template_path.exists():
            raise FileNotFoundError(f"PPT模板文件不存在: {self.template_path}")
            
        if not self.excel_path.exists():
            raise FileNotFoundError(f"Excel文件不存在: {self.excel_path}")
            
        if not self.images_path.exists():
            raise FileNotFoundError(f"图片文件夹不存在: {self.images_path}")
            
        logger.info("所有文件路径验证通过")
    
    def update_first_slide_date(self, presentation):
        """
        更新第一页的日期
        
        Args:
            presentation: PPT演示文稿对象
        """
        logger.info("更新第一页日期...")
        
        try:
            first_slide = presentation.slides[0]
            current_date = datetime.now().strftime("%Y%m%d")
            
            # 查找包含日期的文本框
            for shape in first_slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    # 查找日期格式 (8位数字)
                    date_pattern = r'\d{8}'
                    if re.search(date_pattern, text):
                        # 替换日期
                        new_text = re.sub(date_pattern, current_date, text)
                        shape.text_frame.text = new_text
                        logger.info(f"日期已更新: {text} -> {new_text}")
                        break
            else:
                logger.warning("未找到日期占位符")
                
        except Exception as e:
            logger.error(f"更新日期时发生错误: {e}")
    
    def read_excel_data(self):
        """
        读取Excel数据
        
        Returns:
            pandas.DataFrame: Excel数据
        """
        logger.info("读取Excel数据...")
        
        try:
            # 读取Excel文件
            df = pd.read_excel(self.excel_path, sheet_name="Gemba巡厂 V2")
            
            # 只保留需要的列
            required_columns = ["问题发现区域", "发现人", "问题收集", "问题分类"]
            
            # 检查列是否存在
            for col in required_columns:
                if col not in df.columns:
                    logger.warning(f"Excel中缺少列: {col}")
            
            # 过滤空行
            df = df.dropna(subset=["问题收集"])
            
            logger.info(f"成功读取 {len(df)} 行数据")
            return df
            
        except Exception as e:
            logger.error(f"读取Excel数据时发生错误: {e}")
            raise
    
    def find_matching_image(self, problem_description):
        """
        根据问题描述查找匹配的图片
        
        Args:
            problem_description (str): 问题描述
            
        Returns:
            Path or None: 匹配的图片路径
        """
        if not problem_description or pd.isna(problem_description):
            return None
            
        # 清理问题描述，移除特殊字符
        clean_description = str(problem_description).strip()
        
        # 在图片文件夹中查找匹配的图片
        for image_file in self.images_path.glob("*.jpeg"):
            image_name = image_file.stem  # 不包含扩展名的文件名
            
            # 检查问题描述是否包含在图片文件名中
            if clean_description in image_name:
                logger.info(f"找到匹配图片: {clean_description} -> {image_file.name}")
                return image_file
        
        # 如果没有精确匹配，尝试部分匹配
        for image_file in self.images_path.glob("*.jpeg"):
            image_name = image_file.stem
            
            # 检查图片文件名是否包含在问题描述中
            if image_name in clean_description:
                logger.info(f"找到部分匹配图片: {clean_description} -> {image_file.name}")
                return image_file
        
        logger.warning(f"未找到匹配的图片: {clean_description}")
        return None
    
    def update_category_options(self, slide, category):
        """
        在PPT幻灯片中根据分类更新选项打勾
        
        Args:
            slide: PPT幻灯片对象
            category (str): 问题分类
        """
        if not category or pd.isna(category):
            return
            
        category = str(category).strip()
        logger.info(f"为分类添加勾选标记: {category}")
        
        try:
            # 查找包含分类选项的文本框
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    
                    # 检查是否包含目标分类
                    if category in text:
                        # 在分类文字后添加√符号
                        if "√" not in text:
                            new_text = text.replace(category, f"{category} √")
                            shape.text_frame.text = new_text
                            logger.info(f"已添加勾选标记: {text} -> {new_text}")
                            break
            
        except Exception as e:
            logger.error(f"更新分类选项时发生错误: {e}")
    
    def fill_placeholders(self, slide, area, person, problem):
        """
        填充幻灯片占位符
        
        Args:
            slide: PPT幻灯片对象
            area (str): 问题发现区域
            person (str): 发现人
            problem (str): 问题收集
        """
        logger.info("填充占位符数据...")
        
        try:
            placeholder_map = {
                "模具": str(area) if not pd.isna(area) else "",
                "-": str(person) if not pd.isna(person) else "",
                "看板信息更新": str(problem) if not pd.isna(problem) else ""
            }
            
            # 查找并替换占位符
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    current_text = shape.text_frame.text
                    
                    # 检查是否是需要替换的占位符
                    for placeholder, new_value in placeholder_map.items():
                        if current_text.strip() == placeholder:
                            shape.text_frame.text = new_value
                            logger.info(f"占位符已更新: {placeholder} -> {new_value}")
                            break
            
        except Exception as e:
            logger.error(f"填充占位符时发生错误: {e}")
    
    def add_image_to_slide(self, slide, image_path):
        """
        向幻灯片添加图片
        
        Args:
            slide: PPT幻灯片对象
            image_path (Path): 图片路径
        """
        if not image_path or not image_path.exists():
            return
            
        try:
            # 查找图片占位符
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # 图片占位符
                    # 删除现有图片占位符
                    slide.shapes._spTree.remove(shape._element)
                    break
            
            # 添加新图片
            left = Inches(6)  # 图片位置
            top = Inches(2)
            width = Inches(3)
            height = Inches(2.5)
            
            slide.shapes.add_picture(str(image_path), left, top, width, height)
            logger.info(f"图片已添加到幻灯片: {image_path.name}")
            
        except Exception as e:
            logger.error(f"添加图片时发生错误: {e}")
    
    def generate_ppt(self):
        """
        生成完整的PPT文件
        
        Returns:
            str: 输出文件路径
        """
        logger.info("开始生成PPT文件...")
        
        try:
            # 读取PPT模板
            presentation = Presentation(str(self.template_path))
            logger.info(f"已加载PPT模板: {self.template_path}")
            
            # 更新第一页日期
            self.update_first_slide_date(presentation)
            
            # 读取Excel数据
            df = self.read_excel_data()
            
            # 获取模板幻灯片（假设第二张幻灯片是模板）
            if len(presentation.slides) < 2:
                logger.error("PPT模板至少需要2张幻灯片（第1张为首页，第2张为数据模板）")
                raise ValueError("PPT模板格式不正确")
                
            template_slide = presentation.slides[1]
            
            # 为每行数据创建新幻灯片
            for index, row in df.iterrows():
                logger.info(f"处理第 {index + 1}/{len(df)} 行数据...")
                
                # 复制模板幻灯片
                slide_layout = template_slide.slide_layout
                new_slide = presentation.slides.add_slide(slide_layout)
                
                # 复制模板幻灯片的所有形状
                for shape in template_slide.shapes:
                    try:
                        # 这里需要更复杂的形状复制逻辑
                        # 简化处理：只复制文本框
                        if hasattr(shape, "text_frame"):
                            # 在新幻灯片中创建相同的文本框
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height
                            
                            new_textbox = new_slide.shapes.add_textbox(left, top, width, height)
                            new_textbox.text_frame.text = shape.text_frame.text
                    except Exception as e:
                        logger.warning(f"复制形状时发生错误: {e}")
                
                # 填充占位符
                self.fill_placeholders(
                    new_slide,
                    row.get("问题发现区域"),
                    row.get("发现人"),
                    row.get("问题收集")
                )
                
                # 更新分类选项
                self.update_category_options(new_slide, row.get("问题分类"))
                
                # 添加匹配的图片
                image_path = self.find_matching_image(row.get("问题收集"))
                if image_path:
                    self.add_image_to_slide(new_slide, image_path)
            
            # 删除原始模板幻灯片
            if len(presentation.slides) > 2:
                slides = list(presentation.slides)
                presentation.slides._sldIdLst.remove(slides[1]._element)
            
            # 保存生成的PPT
            output_filename = f"Gemba巡厂报告_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            output_path = self.base_path / output_filename
            
            presentation.save(str(output_path))
            logger.info(f"PPT文件已生成: {output_path}")
            
            return str(output_path)
            
        except Exception as e:
            logger.error(f"生成PPT时发生错误: {e}")
            raise

def main():
    """主函数"""
    try:
        # 设置基础路径
        base_path = r"C:\Users\86151\Desktop\巡厂自动PPT"
        
        # 创建生成器实例
        generator = GembaPPTGenerator(base_path)
        
        # 生成PPT
        output_file = generator.generate_ppt()
        
        print(f"\n✅ PPT生成成功!")
        print(f"📁 输出文件: {output_file}")
        
    except Exception as e:
        logger.error(f"程序执行失败: {e}")
        print(f"\n❌ 程序执行失败: {e}")
        return 1
    
    return 0

if __name__ == "__main__":
    exit_code = main()
    sys.exit(exit_code)