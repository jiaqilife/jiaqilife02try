#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
真正的PPT生成器 - 创建多页幻灯片
"""

import os
import sys
from datetime import datetime
from pathlib import Path
import re

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("python-pptx库未安装，请运行: pip install python-pptx")

def read_excel_data_simple():
    """简化版Excel数据读取 - 使用模拟数据"""
    # 基于实际Excel数据的模拟
    data = [
        {"问题发现区域": "包装", "发现人": "谢佳", "问题收集": "码垛机器人旁边漏雨", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "成品库虚线还要有", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "板台里放了箱子，要分开", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "主路不放木箱", "问题分类": "5S"},
        {"问题发现区域": "成品库、空柄库", "发现人": "谢佳", "问题收集": "这个区域少放料", "问题分类": "5S"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "餐厅区域，信息公布栏，过期信息", "问题分类": "Others"},
        {"问题发现区域": "公共区域", "发现人": "谢佳", "问题收集": "二期餐厅外面空调挂机铁板锈严重", "问题分类": "5S"},
        {"问题发现区域": "装配", "发现人": "谢佳", "问题收集": "立牌子，调试中", "问题分类": "5S"},
        {"问题发现区域": "装配", "发现人": "谢佳", "问题收集": "UV贴纸区域，无关物料不能放在现场", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "办公室上面的玻璃要擦", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "油，管子，清理 刷漆", "问题分类": "5S"},
        {"问题发现区域": "钳子冷锻", "发现人": "谢佳", "问题收集": "玻璃需要擦", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "AGV会看该区域", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "漏雨，电镀门口", "问题分类": "5S"},
        {"问题发现区域": "电镀", "发现人": "谢佳", "问题收集": "下雨，水帘洞", "问题分类": "5S"},
    ]
    
    print(f"读取到 {len(data)} 行数据")
    return data

def update_first_slide_date(presentation):
    """更新第一张幻灯片的日期"""
    try:
        first_slide = presentation.slides[0]
        current_date = datetime.now().strftime("%Y%m%d")
        
        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text
                # 查找8位数字的日期格式
                date_pattern = r'\d{8}'
                if re.search(date_pattern, text):
                    new_text = re.sub(date_pattern, current_date, text)
                    shape.text_frame.text = new_text
                    print(f"日期已更新: {text} -> {new_text}")
                    break
        else:
            print("未找到日期占位符")
            
    except Exception as e:
        print(f"更新日期时发生错误: {e}")

def find_matching_image(problem_description, images_path):
    """查找匹配的图片"""
    if not problem_description or not images_path.exists():
        return None
        
    # 精确匹配
    for image_file in images_path.glob("*.jpeg"):
        image_name = image_file.stem
        if problem_description in image_name:
            return image_file
    
    # 部分匹配 - 检查图片名是否在问题描述中
    for image_file in images_path.glob("*.jpeg"):
        image_name = image_file.stem
        if image_name in problem_description:
            return image_file
    
    return None

def fill_slide_content(slide, area, person, problem, category):
    """填充幻灯片内容"""
    try:
        # 定义占位符映射
        placeholder_map = {
            "模具": area if area else "",
            "-": person if person else "",
            "看板信息更新": problem if problem else ""
        }
        
        # 查找并替换占位符
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                current_text = shape.text_frame.text.strip()
                
                # 替换占位符内容
                for placeholder, new_value in placeholder_map.items():
                    if current_text == placeholder:
                        shape.text_frame.text = new_value
                        print(f"  占位符更新: {placeholder} -> {new_value}")
                        break
                
                # 处理分类打勾
                if category and category in current_text:
                    if "√" not in current_text:
                        new_text = current_text.replace(category, f"{category} √")
                        shape.text_frame.text = new_text
                        print(f"  分类打勾: {current_text} -> {new_text}")
                        
    except Exception as e:
        print(f"填充幻灯片内容时发生错误: {e}")

def add_image_to_slide(slide, image_path):
    """向幻灯片添加图片"""
    if not image_path or not image_path.exists():
        return
        
    try:
        # 在幻灯片右侧添加图片
        left = Inches(6)
        top = Inches(1.5)
        width = Inches(3.5)
        height = Inches(3)
        
        slide.shapes.add_picture(str(image_path), left, top, width, height)
        print(f"  图片已添加: {image_path.name}")
        
    except Exception as e:
        print(f"添加图片时发生错误: {e}")

def create_real_ppt():
    """创建真正的多页PPT"""
    if not PPTX_AVAILABLE:
        print("错误: python-pptx库未安装")
        return None
    
    print("开始创建真正的多页PPT...")
    
    # 设置路径
    base_path = Path(r"C:\Users\86151\Desktop\巡厂自动PPT")
    template_path = base_path / "参观路线Gemba20250829.pptx"
    images_path = base_path / "Gemba巡厂_V2_20250920170854" / "Files" / "待整改--现场图片"
    
    # 检查模板文件
    if not template_path.exists():
        print(f"错误: PPT模板文件不存在: {template_path}")
        return None
    
    try:
        # 加载PPT模板
        presentation = Presentation(str(template_path))
        print(f"已加载PPT模板，当前有 {len(presentation.slides)} 张幻灯片")
        
        # 更新第一张幻灯片的日期
        update_first_slide_date(presentation)
        
        # 获取数据
        data = read_excel_data_simple()
        
        # 获取模板幻灯片（假设第2张是数据模板）
        if len(presentation.slides) < 2:
            print("警告: PPT模板只有一张幻灯片，将使用第一张作为模板")
            template_slide = presentation.slides[0]
        else:
            template_slide = presentation.slides[1]
        
        print(f"\n开始为 {len(data)} 行数据创建幻灯片...")
        
        # 为每行数据创建新幻灯片
        for i, row in enumerate(data, 1):
            print(f"\n--- 处理第 {i}/{len(data)} 行数据 ---")
            print(f"区域: {row['问题发现区域']}")
            print(f"发现人: {row['发现人']}")
            print(f"问题: {row['问题收集']}")
            print(f"分类: {row['问题分类']}")
            
            # 复制模板幻灯片的布局
            slide_layout = template_slide.slide_layout
            new_slide = presentation.slides.add_slide(slide_layout)
            
            # 复制模板幻灯片的内容
            for shape in template_slide.shapes:
                try:
                    if hasattr(shape, "text_frame"):
                        # 创建对应的文本框
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        new_textbox = new_slide.shapes.add_textbox(left, top, width, height)
                        new_textbox.text_frame.text = shape.text_frame.text
                        
                        # 复制文本格式
                        if shape.text_frame.paragraphs:
                            for j, para in enumerate(shape.text_frame.paragraphs):
                                if j < len(new_textbox.text_frame.paragraphs):
                                    new_para = new_textbox.text_frame.paragraphs[j]
                                    if para.runs:
                                        for run in para.runs:
                                            new_para.text = run.text
                                            break
                except Exception as e:
                    print(f"  复制形状时发生警告: {e}")
            
            # 填充幻灯片内容
            fill_slide_content(
                new_slide,
                row['问题发现区域'],
                row['发现人'],
                row['问题收集'],
                row['问题分类']
            )
            
            # 查找并添加匹配的图片
            image_path = find_matching_image(row['问题收集'], images_path)
            if image_path:
                add_image_to_slide(new_slide, image_path)
                print(f"  匹配图片: {image_path.name}")
            else:
                print(f"  未找到匹配图片")
        
        # 删除原始模板幻灯片（保留第一张）
        if len(presentation.slides) > len(data) + 1:
            slides_to_remove = []
            for i in range(1, len(presentation.slides) - len(data)):
                slides_to_remove.append(presentation.slides[i])
            
            for slide in slides_to_remove:
                slide_id = slide.slide_id
                presentation.part.drop_rel(slide.part.partname)
                for slide_rel in presentation.slides._sldIdLst:
                    if slide_rel.id == slide_id:
                        presentation.slides._sldIdLst.remove(slide_rel)
                        break
        
        # 保存生成的PPT
        output_filename = f"Gemba巡厂多页报告_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        output_path = base_path / output_filename
        
        presentation.save(str(output_path))
        
        print(f"\n✅ 多页PPT已生成!")
        print(f"📁 文件位置: {output_path}")
        print(f"📊 总页数: {len(presentation.slides)} 页（首页 + {len(data)} 页数据）")
        
        return str(output_path)
        
    except Exception as e:
        print(f"生成PPT时发生错误: {e}")
        return None

def main():
    """主函数"""
    try:
        output_file = create_real_ppt()
        
        if output_file:
            print("\n程序执行完成!")
            return 0
        else:
            print("\n程序执行失败")
            return 1
            
    except Exception as e:
        print(f"\n程序执行失败: {e}")
        return 1

if __name__ == "__main__":
    exit_code = main()
    sys.exit(exit_code)